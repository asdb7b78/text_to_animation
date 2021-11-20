import sys
import os
import argparse
from pptx import Presentation
from pptx.util import Pt
import imageio

if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('pptx', action='store', help='Pptx file.')
    parser.add_argument('page', action='store',type=int, help='Page number.')
    parser.add_argument('-s', '--sound', action='store', help='Wav or mp3 file.')
    parser.add_argument('-g', '--gif', action='store', help='Gif file.')
    parser.add_argument('--size', action='store_true', help='Resize the gif.')
    parser.add_argument('--height', action='store_true', help='Resize the gif by height.')
    parser.add_argument('--pos', action='store_true', help='Set gif position, default = (0, 0).')
       
    args = parser.parse_args()
    
    prs = Presentation(args.pptx)

    pos = [0, 0]
    size = [0, 0]
    

    if args.gif != None:
        if args.pos:
            print('X Y:')
            pos = input().split(' ')
        if args.size:
            print('Width Height:')
            size = input().split(' ')
        else:
            gif = imageio.imread(args.gif)
            size[0] = gif.shape[0]
            size[1] = gif.shape[1]

        prs.slides[args.page-1].shapes.add_movie(args.gif, Pt(pos[0]), Pt(pos[1]), Pt(size[0]), Pt(size[1]))
    
    if args.sound != None:
        prs.slides[args.page-1].shapes.add_movie(args.sound, Pt(0), Pt(0), Pt(100), Pt(100))

    prs.save('new_'+args.pptx)
