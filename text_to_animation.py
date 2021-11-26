import os
import imageio
from pptx import Presentation
from pptx.util import Pt
from google.cloud import texttospeech

def pptx_to_text(pptx_path, page):
    prs = Presentation(pptx_path)
    text = ''
    for shape in prs.slides[page-1].shapes:
        if shape.has_text_frame and (shape.text != ''):
            print(shape.text+":(y/n)")
            while(True):
                tmp_in = input()
                if tmp_in == 'y':
                    text += shape.text + ' '
                    break
                elif tmp_in == 'n':
                    break
                else:
                    print("Please input y or n")
    print("\nAll texts to trans: " + text)
    return text

def text_to_speech(text_in, wav_path):
    """
    Synthesizes speech from the input string of text or ssml.
    Note: ssml must be well-formed according to:
        https://www.w3.org/TR/speech-synthesis/
    """
    # Instantiates a client
    client = texttospeech.TextToSpeechClient()
    # Set the text input to be synthesized
    synthesis_input = texttospeech.SynthesisInput(text=text_in)
    # Build the voice request, select the language code ("en-US") and the ssml
    # voice gender ("neutral")
    voice = texttospeech.VoiceSelectionParams(
        language_code="en-US", ssml_gender=texttospeech.SsmlVoiceGender.NEUTRAL
    )
    # Select the type of audio file you want returned
    audio_config = texttospeech.AudioConfig(
        audio_encoding=texttospeech.AudioEncoding.LINEAR16
    )
    # Perform the text-to-speech request on the text input with the selected
    # voice parameters and audio file type
    response = client.synthesize_speech(
        input=synthesis_input, voice=voice, audio_config=audio_config
    )
    # The response's audio_content is binary.
    with open(wav_path, "wb") as out:
        # Write the response to the output file.
        out.write(response.audio_content)
        print('\nAudio content written to file "'+ wav_path + '\n')

def make_gif(folder, gif_parh, dur, mouth_path):
    # Read PNG
    img_name = ['A', 'B', 'C', 'D', 'E', 'F', 'G', 'H']
    img = []
    for n in img_name:
        img.append(imageio.imread('./Png/'+folder+'/'+n+'.png'))

    # Read mouth.txt
    mouth = []
    f = open(mouth_path)
    for line in f:
        i = line.split('\t')
        mouth.append([i[0], i[1][0]])

    for i in range(0, len(mouth)-1):
        mouth[i][0] = float(mouth[i+1][0]) - float(mouth[i][0])
    mouth.pop(len(mouth)-1)

    # Generate gif list
    gif = []
    for i in mouth:
        # Round the duration
        i[0] = int((float(i[0]) + dur/2)//dur)
        for j in range(0, i[0]):
            if i[1] == 'X' or i[1] == 'A':
                gif.append(img[0])
            elif i[1] == 'B':
                gif.append(img[1])
            elif i[1] == 'C':
                gif.append(img[2])
            elif i[1] == 'D':
                gif.append(img[3])
            elif i[1] == 'E':
                gif.append(img[4])
            elif i[1] == 'F':
                gif.append(img[5])
            elif i[1] == 'G':
                gif.append(img[6])
            elif i[1] == 'H':
                gif.append(img[7])
    imageio.mimsave(gif_parh, gif, duration = dur, loop=1)
    print('\nGif is written to file "' + gif_parh + '"\n')

def insert_to_pptx(pptx_path, page, gif_path, sound_path, pos, size, height):
    prs = Presentation(pptx_path)
    if pos == False:
        pos = [0, 0]
    if size == False:
        size = [0, 0]
        gif = imageio.imread(gif_path)
        size[0] = gif.shape[0]
        size[1] = gif.shape[1]
        if height != False:
            size[0] = int(float(size[0]) * height / size[1])
            size[1] = height
    prs.slides[page-1].shapes.add_picture(gif_path, Pt(pos[0]), Pt(pos[1]), width=Pt(size[0]), height=Pt(size[1]))    
    prs.slides[page-1].shapes.add_movie(sound_path, Pt(0), Pt(0), Pt(100), Pt(100))
    return prs

def save_pptx(prs, path):
    prs.save(path)
    print('\nNew pptx is written to file "'+path+'"\n')

def text_to_animation(pptx, page, name, folder, duration, pos=False, size=False, height=False):
    # pptx = path of pptx
    # page = page numbeer
    # name = output gif, wav name
    # folder = source image folder in 'PNG'
    # duration = minimum duration of every source image in gif
    # pos = [x, y] is position of gif,
    #       if set True, get a new pos from standard input
    # size = [w, h] is size of gif,
    #       if set True, get new w and h from standard input
    # height = h is height of gif, and would scale up the gif by h,
    #       if set True, get new h from standard input

    text = pptx_to_text(pptx, page)
    text_to_speech(text, name+'.wav')

    mouth = 'mouth.txt'
    os.system("rhubarb.exe " + name + ".wav -o "+mouth)
    make_gif(folder, name + '.gif', dur=duration, mouth_path=mouth)

    if pos == True:
        print('Input new pos X Y in one line:')
        pos = input().split(' ')
        pos[0] = int(pos[0])
        pos[1] = int(pos[1])
    if size != False and height != False:
        print("Can't resize gif with size and height simultaneously.")
        return
    if size == True:
        print('Input new width and height in one line:')
        size = input().split(' ')
        size[0] = int(size[0])
        size[1] = int(size[1])
    elif height == True:
        print('Input new height:')
        height = int(input())

    prs = insert_to_pptx(pptx, page, gif_path=name+'.gif', sound_path=name+'.wav', pos=pos, size=size, height=height)
    save_pptx(prs, path='new_'+pptx)
