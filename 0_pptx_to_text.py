import sys
import os
import argparse
from pptx import Presentation

#讀檔 -> 指定哪一頁 -> 轉換文字 -> 插入gif -> 插入聲音
if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('pptx', action='store', help='Pptx name')
    parser.add_argument('page', action='store', default='0', type=int, help='Page number')
    args = parser.parse_args()

    prs = Presentation(args.pptx)

    text = ''
    for shape in prs.slides[args.page-1].shapes:
        if shape.has_text_frame and (shape.text != ''):
            print(shape.text)
            if input() == 'y':
                text += shape.text+' '
    print(text)