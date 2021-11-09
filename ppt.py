# -*- coding: UTF-8 -*-
from pptx import Presentation
#開啟新的簡報物件
prs = Presentation()
#建立簡報檔第一張頁面物件
title_slide_layout = prs.slide_layouts[0] 
#增加一張簡報
slide = prs.slides.add_slide(title_slide_layout)
#設定第一張簡報的標題 
title = slide.shapes.title
title.text = "Hello Python PPT"
#設定第一張簡報的副標題
#subtitle = slide.placeholders[1]
#subtitle.text = "作者：Meiko 2020/10/01"

#insert gif
pic = slide.shapes.add_picture('hello.gif', 15, 15)


#將簡報物件存檔
prs.save("python_ppt.pptx")
