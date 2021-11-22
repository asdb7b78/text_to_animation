import sys
import os
import time
import argparse
import imageio
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import PP_MEDIA_TYPE
from google.cloud import texttospeech


def text_to_speech(text_in, wav_name):
    """
    Synthesizes speech from the input string of text or ssml.
    Note: ssml must be well-formed according to:
        https://www.w3.org/TR/speech-synthesis/
    """
    # Instantiates a client
    client = texttospeech.TextToSpeechClient()
    # Set the text input to be synthesized
    synthesis_input = texttospeech.SynthesisInput(text=text_in)
    # Build the voice request, select the language code ("en-US") and the ssml
    # voice gender ("neutral")
    voice = texttospeech.VoiceSelectionParams(
        language_code="en-US", ssml_gender=texttospeech.SsmlVoiceGender.NEUTRAL
    )
    # Select the type of audio file you want returned
    audio_config = texttospeech.AudioConfig(
        audio_encoding=texttospeech.AudioEncoding.LINEAR16
    )
    # Perform the text-to-speech request on the text input with the selected
    # voice parameters and audio file type
    response = client.synthesize_speech(
        input=synthesis_input, voice=voice, audio_config=audio_config
    )
    # The response's audio_content is binary.
    with open(args.name + ".wav", "wb") as out:
        # Write the response to the output file.
        out.write(response.audio_content)
        print('\nAudio content written to file "'+ args.name +'.wav"\n')

def make_gif(folder, gif_name, dur):
    
    # Read PNG
    img_name = ['A', 'B', 'C', 'D', 'E', 'F', 'G', 'H']
    img = []
    for n in img_name:
        img.append(imageio.imread('./Png/'+folder+'/'+n+'.png'))
    
    # Read mouth.txt
    mouth = []
    f = open("mouth.txt")
    for line in f:
        i = line.split('\t')
        mouth.append([i[0], i[1][0]])

    for i in range(0, len(mouth)-1):
        mouth[i][0] = float(mouth[i+1][0]) - float(mouth[i][0])
    mouth.pop(len(mouth)-1)

    # Generate gif list
    gif = []
    for i in mouth:
        # Round the duration
        i[0] = int((float(i[0]) + dur/2)//dur)
        for j in range(0, i[0]):
            if i[1] == 'X' or i[1] == 'A':
                gif.append(img[0])
            elif i[1] == 'B':
                gif.append(img[1])
            elif i[1] == 'C':
                gif.append(img[2])
            elif i[1] == 'D':
                gif.append(img[3])
            elif i[1] == 'E':
                gif.append(img[4])
            elif i[1] == 'F':
                gif.append(img[5])
            elif i[1] == 'G':
                gif.append(img[6])
            elif i[1] == 'H':
                gif.append(img[7])
    imageio.mimsave(gif_name+'.gif', gif, duration = dur, loop=1)
    print('\nGif is written to file "' + gif_name + '.gif"\n')

#讀檔 -> 指定哪一頁 -> 轉換文字 -> 插入gif -> 插入聲音
if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('pptx', action='store', help='Pptx name')
    parser.add_argument('page', action='store', default='0', type=int, help='Page number')
    parser.add_argument('-n', '--name', action='store', default='output', help="Name of .wav and .gif")
    parser.add_argument('-f', '--folder', action='store',default='Lip', help='Folder name in "PNG"')
    parser.add_argument('-d', '--duration', action='store',default='0.2', type=float, help='Gif duration of each image, default=0.2')
    parser.add_argument('--size', action='store_true', help='Resize the gif.')
    parser.add_argument('--height', action='store_true', help='Resize the gif by height.')
    parser.add_argument('--pos', action='store_true', help='Set gif position, default = (0, 0).')
    args = parser.parse_args()

    prs = Presentation(args.pptx)

    text = ''
    for shape in prs.slides[args.page-1].shapes:
        if shape.has_text_frame and (shape.text != ''):
            print(shape.text+":(y/n)")
            tmp_in = input()
            if tmp_in == 'y':
                text += shape.text + ' '
            elif tmp_in == 'n':
                continue
    print("\nAll texts are: " + text)

    text_to_speech(text, args.name)
    
    os.system("./rhubarb " + args.name + ".wav -o mouth.txt")
    
    make_gif(args.folder, args.name, args.duration)

    pos = [0, 0]
    size = [300, 300]

    if args.pos:
        print('X Y:')
        pos = input().split(' ')
        pos[0] = int(pos[0])
        pos[1] = int(pos[1])
    if args.size:
        print('Width Height:')
        size = input().split(' ')
        size[0] = int(size[0])
        size[1] = int(size[1])
    else:
        gif = imageio.imread(args.name+'.gif')
        size[0] = gif.shape[0]
        size[1] = gif.shape[1]
        if args.height:
            print('Height:')
            h = int(input())
            size[0] = int(float(size[0]) * h / size[1])
            size[1] = h
            

    #prs.slides[args.page-1].shapes.add_movie(args.name + '.gif', Pt(pos[0]), Pt(pos[1]), Pt(size[0]), Pt(size[1]))
    prs.slides[args.page-1].shapes.add_picture(args.name + '.gif', Pt(pos[0]), Pt(pos[1]), width=Pt(size[0]), height=Pt(size[1]))    
    prs.slides[args.page-1].shapes.add_movie(args.name + '.wav', Pt(0), Pt(0), Pt(100), Pt(100))

    prs.save('new_' + args.pptx)
    print('\nNew pptx is written to file "new_' + args.pptx + '.gif"\n')
